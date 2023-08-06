# coding: utf-8
import os
import warnings

import pandas as pd
import docx
from pptx import Presentation as PowerPoint
from pptx.util import Cm, Pt

from .. import PathExistsWarning


def Excel(io, sheetname=0, header=0):
    frames = pd.read_excel(io, sheetname=sheetname, header=header)
    if isinstance(frames, dict):
        frames = [(name, frame) for name, frame in frames.items()]
        frames.sort()
    return frames


def to_excel(src, path, overwrite=False):
    """Helper to save pandas.DataFrame(s) to excel file.

    pandas.DataFrame.to_excel requires steps
    to save multiple frames as a single file.
    This helper function saves multiple frames
    into a single excel file in a single call.

    :param
    src: DataFrame or list of DataFrames
    path: excel file path to save DataFrame(s)
    """
    if not path.endswith('.xlsx'):
        path = '{}.xlsx'.format(path)

    if os.path.exists(path) and not overwrite:
        warnings.warn(PathExistsWarning(path))

    if isinstance(src, pd.DataFrame):
        return src.to_excel(path)

    excel_writer = pd.ExcelWriter(path)
    for i, frame in enumerate(src):
        if hasattr(frame, 'name'):
            sheetname = frame.name
        else:
            sheetname = 'sheet {}'.format(i)
        frame.to_excel(excel_writer, sheetname)
    else:
        excel_writer.save()


def Word(filepath_or_buffer=None):
    """
    Extends docx.document.Document
    """
    # patch class
    docx.document.Document.extract_text = extract_text
    docx.document.Document.extract_tables = extract_tables
    docx.document.Document.tables_to_excel = tables_to_excel
    docx.document.Document.add_tables_from_excel = add_tables_from_excel
    docx.document.Document.add_picture = add_picture

    # make instance using patched class
    return docx.Document(filepath_or_buffer)


def extract_text(self):
    return '\n'.join([p.text for p in self.paragraphs])


def extract_tables(self, select=None):
    table_list = []
    for table in self.tables:
        table_data = []
        for row in table.rows:
            row_data = []
            for cell in row.cells:
                cell_text = ''
                for p in cell.paragraphs:
                    cell_text += p.text
                # if debug:
                #     print(cell_text, end=' ')
                row_data.append(cell_text)
            # if debug:
            #     print()
            table_data.append(row_data)

        table_frame = pd.DataFrame(table_data)
        table_list.append(table_frame)

    if select is None:
        return table_list
    else:
        return table_list[select]


def tables_to_excel(self, filepath_or_buffer):
    table_list = self.extract_tables()

    excel_file = pd.ExcelWriter(filepath_or_buffer)
    for i, table_frame in enumerate(table_list):
        table_frame.to_excel(excel_file, 'Sheet{}'.format(i+1))
    excel_file.save()

    return filepath_or_buffer


def add_tables_from_excel(
        self, excel_filepath, sheetname=None, table_style='TableGrid'):
    table_frames = Excel(excel_filepath, sheetname=sheetname)

    if not isinstance(table_frames, list):
        table_frames = [table_frames]

    for frame in table_frames:
        frame = frame.fillna('')
        nrows, ncols = len(frame), len(frame.columns)
        table = self.add_table(rows=nrows, cols=ncols, style='TableGrid')
        # table.style = 'LightShading-Accent1'
        for i, row in enumerate(table.rows):
            for j, cell in enumerate(row.cells):
                cell.text = str(frame.iloc[i, j])


def add_picture(self, image_path_or_stream, width=None, height=None):
    if width:
        width = docx.shared.Inches(width)
    if height:
        height = docx.shared.Inches(height)

    run = self.add_paragraph().add_run()
    return run.add_picture(image_path_or_stream, width, height)
